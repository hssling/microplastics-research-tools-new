from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
import matplotlib.pyplot as plt

# Load the existing presentation
pptx_path = "/mnt/data/Epidemiology_Basics_Presentation.pptx"
prs = Presentation(pptx_path)

# Function to add an image slide
def add_image_slide(prs, title, image_path):
    slide_layout = prs.slide_layouts[5]  # blank layout
    slide = prs.slides.add_slide(slide_layout)
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(8.5)
    height = Inches(6)
    slide.shapes.add_picture(image_path, left, top, width=width, height=height)
    # Add title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0), Inches(9), Inches(0.5))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Inches(0.4)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

# --- Create and save illustrative diagrams ---

# 1. Epidemiological Triad
plt.figure(figsize=(5,5))
plt.text(0.5, 0.9, "Agent", ha="center", fontsize=14, bbox=dict(facecolor="lightblue"))
plt.text(0.1, 0.1, "Host", ha="center", fontsize=14, bbox=dict(facecolor="lightgreen"))
plt.text(0.9, 0.1, "Environment", ha="center", fontsize=14, bbox=dict(facecolor="lightyellow"))
plt.plot([0.5,0.1],[0.85,0.15],'k-')
plt.plot([0.5,0.9],[0.85,0.15],'k-')
plt.plot([0.1,0.9],[0.15,0.15],'k-')
plt.axis("off")
triad_img = "/mnt/data/epidemiological_triad.png"
plt.savefig(triad_img, bbox_inches="tight")
plt.close()

add_image_slide(prs, "Epidemiological Triad", triad_img)

# 2. Incidence vs Prevalence
plt.figure(figsize=(6,4))
time = list(range(10))
incidence = [1,2,1,3,2,2,1,2,1,1]
prevalence = [sum(incidence[:i+1]) for i in range(len(incidence))]
plt.plot(time, incidence, label="Incidence (new cases)", marker="o")
plt.plot(time, prevalence, label="Prevalence (all cases)", marker="s")
plt.xlabel("Time")
plt.ylabel("Cases")
plt.title("Incidence vs Prevalence")
plt.legend()
ivp_img = "/mnt/data/incidence_vs_prevalence.png"
plt.savefig(ivp_img, bbox_inches="tight")
plt.close()

add_image_slide(prs, "Incidence vs Prevalence", ivp_img)

# 3. Study Designs Overview (simplified flowchart)
plt.figure(figsize=(6,5))
plt.text(0.5, 0.9, "Epidemiological Studies", ha="center", fontsize=14, bbox=dict(facecolor="lightblue"))
plt.text(0.25, 0.6, "Descriptive", ha="center", fontsize=12, bbox=dict(facecolor="lightgreen"))
plt.text(0.75, 0.6, "Analytic", ha="center", fontsize=12, bbox=dict(facecolor="lightyellow"))
plt.text(0.25, 0.3, "Ecological\nCross-sectional", ha="center", fontsize=10, bbox=dict(facecolor="white"))
plt.text(0.75, 0.3, "Case-control\nCohort", ha="center", fontsize=10, bbox=dict(facecolor="white"))
plt.text(0.5, 0.1, "Experimental (RCTs)", ha="center", fontsize=12, bbox=dict(facecolor="salmon"))
plt.arrow(0.5,0.85, -0.2, -0.2, head_width=0.03, color="black")
plt.arrow(0.5,0.85, 0.2, -0.2, head_width=0.03, color="black")
plt.arrow(0.25,0.55, 0, -0.2, head_width=0.03, color="black")
plt.arrow(0.75,0.55, 0, -0.2, head_width=0.03, color="black")
plt.arrow(0.5,0.85, 0, -0.65, head_width=0.03, color="black")
plt.axis("off")
study_design_img = "/mnt/data/study_designs.png"
plt.savefig(study_design_img, bbox_inches="tight")
plt.close()

add_image_slide(prs, "Study Designs Overview", study_design_img)

# Save updated presentation with graphics
pptx_graphics_path = "/mnt/data/Epidemiology_Basics_with_Illustrations.pptx"
prs.save(pptx_graphics_path)

pptx_graphics_path
