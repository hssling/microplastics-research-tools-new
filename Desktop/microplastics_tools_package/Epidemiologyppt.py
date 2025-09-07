from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Initialize the presentation
prs = Presentation()


# Helper function to add a slide with title, bullets, and optional image
def add_slide_with_bullets(prs, title, bullets, layout_index=1, image_path=None, img_left=8, img_top=2, img_width=2, img_height=2):
    slide_layout = prs.slide_layouts[layout_index]  # 1 = Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    title_placeholder.text = title
    title_placeholder.text_frame.paragraphs[0].font.size = Pt(32)
    title_placeholder.text_frame.paragraphs[0].font.bold = True
    title_placeholder.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()  # Clear default text
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(34, 34, 34)
    # Add image if provided
    if image_path:
        try:
            slide.shapes.add_picture(image_path, Inches(img_left), Inches(img_top), Inches(img_width), Inches(img_height))
        except Exception:
            add_image_placeholder(slide, Inches(img_left), Inches(img_top), Inches(img_width), Inches(img_height))
    return slide


# Helper function to add an image placeholder
def add_image_placeholder(slide, left, top, width, height, text="Image Placeholder"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.text = text
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)


# Slide 1: Title Slide
slide_layout = prs.slide_layouts[0]  # Title Slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Principles of Epidemiology: The Science of Public Health"
title.text_frame.paragraphs[0].font.size = Pt(40)
title.text_frame.paragraphs[0].font.bold = True
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)  # Dark blue

subtitle = slide.placeholders[1]
subtitle.text = "Solving Health Mysteries with Data\n[Your Name/Department]\nChapter 3, Park's Textbook of Preventive and Social Medicine, 27th Ed."
subtitle.text_frame.paragraphs[0].font.size = Pt(22)
subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)
add_image_placeholder(slide, Inches(4), Inches(2), Inches(5), Inches(3), text="Collage: Public Health, Data, Community")



# Slide 2: Learning Objectives
add_slide_with_bullets(prs, "Learning Objectives", [
    "Define epidemiology and its aims with real-world examples (e.g., COVID-19, malaria).",
    "Understand rates, ratios, and proportions (visual: pie chart, bar graph, e.g., disease rates in India vs. USA).",
    "Differentiate incidence vs. prevalence (visual: bathtub analogy, e.g., diabetes prevalence in India).",
    "Explore key epidemiological study designs (icons: cohort, case-control, RCT, e.g., Framingham Heart Study).",
    "Distinguish association vs. causation (Venn diagram, arrows, e.g., smoking and lung cancer).",
    "Outline steps in epidemic investigation (flowchart, e.g., cholera outbreak investigation)."
], image_path="objectives.png", img_left=8, img_top=2, img_width=2, img_height=2)
slide = prs.slides[-1]
add_image_placeholder(slide, Inches(8), Inches(2), Inches(1), Inches(1), text="Icons: Objectives, e.g., magnifying glass, chart")



# Slide 3: What is Epidemiology?
add_slide_with_bullets(prs, "What is Epidemiology?", [
    "Definition: Study of distribution and determinants of health states/events in populations, and application to control health problems.",
    "Example: John Snow's cholera map (London, 1854) – classic epidemiology in action.",
    "Key Questions:",
    "  - Distribution: Who gets the disease? Where? When? (Map, timeline, e.g., COVID-19 spread map)",
    "  - Determinants: Why and how? (Risk factors, exposures, e.g., air pollution and asthma)",
    "  - Application: How do we control/prevent? (Interventions, policies, e.g., vaccination campaigns)"
], image_path="john_snow_map.png", img_left=8, img_top=2, img_width=2, img_height=2)
slide = prs.slides[-1]
add_image_placeholder(slide, Inches(8), Inches(2), Inches(2), Inches(2), text="Flowchart: Distribution → Determinants → Application, e.g., John Snow map")



# Slide 4: Aims of Epidemiology
add_slide_with_bullets(prs, "Aims of Epidemiology", [
    "Goal: Reduce health problems, promote community health (e.g., polio eradication).",
    "Specific Aims:",
    "  1. Describe disease magnitude and distribution (heatmap, bar chart, e.g., malaria incidence map).",
    "  2. Identify risk factors (determinants) (icon: magnifying glass, DNA, e.g., genetic risk for diabetes).",
    "  3. Inform health service planning and evaluation (hospital, checklist, e.g., COVID-19 resource allocation)."
], image_path="aims.png", img_left=8, img_top=2, img_width=1, img_height=1)
slide = prs.slides[-1]
add_image_placeholder(slide, Inches(8), Inches(2), Inches(1), Inches(1), text="Icons: Aims, e.g., hospital, checklist")



# Slide 5: Measuring Health: Rates, Ratios, Proportions
add_slide_with_bullets(prs, "Measuring Health: Rates, Ratios, Proportions", [
    "Rate: Events over time per population (e.g., Crude Death Rate = [Deaths/Mid-year population] × 1000, COVID-19 case rate).",
    "Ratio: Relationship between two quantities (e.g., Male:Female disease ratio, e.g., diabetes ratio in urban vs. rural).",
    "Proportion: Numerator included in denominator, as % (e.g., % of children with scabies, vaccination coverage).",
    "Visuals: Pie chart for proportions, bar graph for rates, ratio icon, e.g., WHO infographics."
], image_path="rates.png", img_left=8, img_top=2, img_width=2, img_height=1)
slide = prs.slides[-1]
add_image_placeholder(slide, Inches(8), Inches(2), Inches(2), Inches(1), text="Diagram: Pie chart, bar graph, ratio, e.g., WHO infographic")



# Slide 6: Incidence vs. Prevalence
add_slide_with_bullets(prs, "Incidence vs. Prevalence", [
    "Incidence: New cases in a population over time (measures risk, e.g., COVID-19 new cases per day).",
    "Prevalence: All existing cases (new + old) at a point/period (measures burden, e.g., diabetes prevalence in India).",
    "Relationship: Prevalence ≈ Incidence × Duration.",
    "Visual: Bathtub analogy (water in tub = prevalence, tap = incidence, drain = recovery/death, e.g., infographic)."
], image_path="bathtub.png", img_left=8, img_top=2, img_width=2, img_height=2)
slide = prs.slides[-1]
add_image_placeholder(slide, Inches(8), Inches(2), Inches(2), Inches(2), text="Bathtub analogy: Incidence, Prevalence, e.g., infographic")



# Slide 7: Epidemiological Study Designs
add_slide_with_bullets(prs, "Epidemiological Study Designs", [
    "Observational: Descriptive (Who, Where, When, e.g., disease mapping), Analytical (Why, How, e.g., risk factor analysis).",
    "Experimental: Tests interventions (e.g., RCTs, vaccine trials, e.g., COVID-19 vaccine studies).",
    "Visuals: Flowchart of study types, icons for each design, e.g., Framingham Heart Study diagram."
], image_path="study_designs.png", img_left=8, img_top=2, img_width=2, img_height=2)
slide = prs.slides[-1]
add_image_placeholder(slide, Inches(8), Inches(2), Inches(2), Inches(2), text="Flowchart: Study Designs, e.g., Framingham Study")



# Slide 8: Case-Control Study
add_slide_with_bullets(prs, "Case-Control Study", [
    "Approach: Retrospective (Effect → Cause, e.g., lung cancer and smoking history).",
    "Process:",
    "  1. Select Cases (with disease) and Controls (without disease).",
    "  2. Compare past exposure to risk factors (e.g., smoking, occupational exposure).",
    "Measure: Odds Ratio (OR), e.g., OR for smoking and lung cancer > 10.",
    "Pros: Quick, cost-effective, ideal for rare diseases (e.g., mesothelioma).",
    "Cons: Prone to bias (e.g., recall bias, selection bias).",
    "Visual: 2 groups, arrows from exposure to outcome, OR formula, e.g., diagram from CDC."
], image_path="case_control.png", img_left=8, img_top=2, img_width=2, img_height=2)
slide = prs.slides[-1]
add_image_placeholder(slide, Inches(8), Inches(2), Inches(2), Inches(2), text="Diagram: Case-Control, e.g., CDC graphic")



# Slide 9: Cohort Study
add_slide_with_bullets(prs, "Cohort Study", [
    "Approach: Prospective (Cause → Effect, e.g., Nurses' Health Study).",
    "Process:",
    "  1. Select disease-free cohort, group by exposure (e.g., smokers vs. non-smokers).",
    "  2. Follow over time to track disease development (e.g., heart disease incidence).",
    "Measures: Relative Risk (RR), Attributable Risk (AR), e.g., RR for smoking and heart disease ≈ 2.5.",
    "Pros: Establishes temporality, calculates incidence, e.g., COVID-19 vaccine effectiveness studies.",
    "Cons: Time-consuming, expensive, not for rare diseases.",
    "Visual: Timeline, exposed vs. unexposed, RR formula, e.g., diagram from NEJM."
], image_path="cohort.png", img_left=8, img_top=2, img_width=2, img_height=2)
slide = prs.slides[-1]
add_image_placeholder(slide, Inches(8), Inches(2), Inches(2), Inches(2), text="Diagram: Cohort Study, e.g., NEJM graphic")



# Slide 10: Randomized Controlled Trial (RCT)
add_slide_with_bullets(prs, "Randomized Controlled Trial (RCT)", [
    "Gold Standard for causality (e.g., COVID-19 vaccine trials, Salk polio vaccine).",
    "Key Features:",
    "  - Intervention: Researcher controls exposure (e.g., drug vs. placebo, e.g., double-blind COVID-19 trial).",
    "  - Randomization: Minimizes selection bias, e.g., random number generator, sealed envelopes.",
    "  - Blinding: Reduces bias (single, double, triple, e.g., placebo-controlled studies).",
    "Visual: Randomization diagram, blinding icons, e.g., CONSORT flowchart."
], image_path="rct.png", img_left=8, img_top=2, img_width=2, img_height=2)
slide = prs.slides[-1]
add_image_placeholder(slide, Inches(8), Inches(2), Inches(2), Inches(2), text="Diagram: RCT, e.g., CONSORT flowchart")



# Slide 11: Association vs. Causation
add_slide_with_bullets(prs, "Association vs. Causation", [
    "Association: Statistical link (may be spurious, indirect, or causal, e.g., ice cream sales and drowning).",
    "Causation: Requires evidence beyond association (e.g., smoking and lung cancer, proven by multiple studies).",
    "Bradford Hill Criteria (select):",
    "  - Temporality: Cause precedes effect (e.g., exposure before disease).",
    "  - Strength: High RR/OR (e.g., RR for smoking and lung cancer > 10).",
    "  - Consistency: Replicated findings (e.g., multiple countries, populations).",
    "  - Biological Plausibility: Makes sense biologically (e.g., carcinogens in tobacco).",
    "Visual: Venn diagram, arrows from association to causation, e.g., infographic from CDC."
], image_path="association_causation.png", img_left=8, img_top=2, img_width=2, img_height=2)
slide = prs.slides[-1]
add_image_placeholder(slide, Inches(8), Inches(2), Inches(2), Inches(2), text="Venn diagram: Association vs. Causation, e.g., CDC infographic")



# Slide 12: Investigating an Epidemic
add_slide_with_bullets(prs, "Investigating an Epidemic", [
    "1. Verify diagnosis and confirm epidemic (e.g., COVID-19, Ebola outbreaks).",
    "2. Define population at risk; find cases (e.g., contact tracing, GIS mapping).",
    "3. Analyze by Time, Place, Person (visual: epidemic curve, map, e.g., COVID-19 dashboard).",
    "4. Formulate and test hypotheses (e.g., contaminated water source).",
    "5. Evaluate ecological factors (e.g., climate, sanitation).",
    "6. Write report (e.g., WHO outbreak reports).",
    "Visual: Flowchart of investigation steps, e.g., CDC outbreak investigation graphic."
], image_path="epidemic_investigation.png", img_left=8, img_top=2, img_width=2, img_height=2)
slide = prs.slides[-1]
add_image_placeholder(slide, Inches(8), Inches(2), Inches(2), Inches(2), text="Flowchart: Epidemic Investigation, e.g., CDC graphic")



# Slide 13: Summary
add_slide_with_bullets(prs, "Summary", [
    "Epidemiology studies disease distribution and determinants (e.g., COVID-19, malaria, diabetes).",
    "Incidence (new cases) and Prevalence (all cases) measure morbidity (e.g., global diabetes prevalence map).",
    "Study designs: Descriptive (patterns), Analytical (hypotheses), Experimental (causality, e.g., vaccine trials).",
    "Case-Control (backward), Cohort (forward), RCTs (controlled, e.g., Salk polio vaccine).",
    "Causation requires rigorous evidence (Bradford Hill, e.g., tobacco and cancer).",
    "Epidemic investigation is systematic and action-oriented (e.g., Ebola outbreak steps).",
    "Visual: Table summarizing study designs and measures, e.g., infographic from WHO."
], image_path="summary_table.png", img_left=8, img_top=2, img_width=2, img_height=1)
slide = prs.slides[-1]
add_image_placeholder(slide, Inches(8), Inches(2), Inches(2), Inches(1), text="Table: Study Designs & Measures, e.g., WHO infographic")


# Slide 14: Thank You
slide_layout = prs.slide_layouts[0]  # Title Slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Thank You"
title.text_frame.paragraphs[0].font.size = Pt(36)
title.text_frame.paragraphs[0].font.bold = True
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
subtitle = slide.placeholders[1]
subtitle.text = "Questions? Contact: [Your Email]"
subtitle.text_frame.paragraphs[0].font.size = Pt(20)
subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)
add_image_placeholder(slide, Inches(4), Inches(2), Inches(5), Inches(3), text="Team/Contact Image")

# Save the presentation
prs.save("Epidemiology_Presentation.pptx")
print("PowerPoint presentation saved as 'Epidemiology_Presentation.pptx'")