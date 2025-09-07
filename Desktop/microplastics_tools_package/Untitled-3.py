### **Instructions to Use the Script**

1.  **Install the necessary library:** Open your terminal or command prompt and run the following command:
    ```bash
    pip install python-pptx
    ```
2.  **Save the code:** Save the Python code below as a file, for example, `create_presentation.py`.
3.  **Prepare your images:** Find the images you want to use for each slide and place them in the same folder as your script. **Update the placeholder image paths** in the script (e.g., change `'image1.jpg'` to the actual name of your file).
4.  **Run the script:** Open your terminal or command prompt, navigate to the folder where you saved the file, and run:
    ```bash
    python create_presentation.py
    ```
5.  A file named `Microplastics_and_Human_Health.pptx` will be created in the same folder. You can then open it and manually add your desired animations.

---

### **Python Script for Presentation Generation**

```python
from pptx import Presentation
from pptx.util import Inches, Pt

# Create a new presentation object
prs = Presentation()

# Set slide width and height for widescreen (16:9)
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# --- Slide 1: Title Slide ---
slide_layout = prs.slide_layouts # Title slide layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders

title.text = "Microplastics and Human Health: A Systematic Review of Exposure, Biodistribution, and Toxicological Outcomes"
subtitle.text = "Author: Dr. Siddalingaiah H S\nProfessor, Department of Community Medicine\nCorrespondence: hssling@yahoo.com"

# --- Slide 2: Introduction ---
slide_layout = prs.slide_layouts # Title and Content layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.shapes.placeholders

title.text = "Introduction: A Pervasive Environmental Issue"
content.text = (
    "• Plastic pollution is a defining environmental issue of the 21st century.\n"
    "• Microplastics (MPs): Synthetic polymer particles smaller than 5 mm.\n"
    "  – Primary MPs: Intentionally manufactured small (e.g., microbeads).\n"
    "  – Secondary MPs: Result from breakdown of larger plastic debris.\n"
    "• Ubiquitous contamination has led to perpetual and unavoidable human exposure.\n"
    "• The issue has transitioned from an ecological concern to a potential public health crisis."
)

# --- Slide 3: Objectives ---
slide_layout = prs.slide_layouts
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.shapes.placeholders

title.text = "Objectives of this Systematic Review"
content.text = (
    "• Primary Goal: To consolidate and critically evaluate current knowledge on the human health implications of microplastics.\n"
    "• Key Focus Areas:\n"
    "  1. Reviewing pathways of human exposure to MPs.\n"
    "  2. Examining their fate within the body (biodistribution).\n"
    "  3. Synthesizing evidence on mechanistic pathways of adverse health effects.\n"
    "  4. Identifying critical knowledge gaps to guide future research."
)

# --- Slide 4: Pathways of Human Exposure ---
slide_layout = prs.slide_layouts
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.shapes.placeholders

title.text = "Pathways of Human Exposure"
content.text = (
    "• Ingestion (Dominant Route):\n"
    "  – Food: Seafood (mussels, oysters), sea salt, honey, sugar, beer.\n"
    "  – Water: Found in tap and bottled water; higher in bottled water.\n"
    "  – Packaging: Abrasion from plastic containers, especially during heating.\n"
    "• Inhalation:\n"
    "  – Sources: Synthetic textiles, car tires, urban dust.\n"
    "  – Indoor Air: Concentrations often higher indoors; fiber-shaped MPs are common.\n"
    "• Dermal Contact:\n"
    "  – Sources: Cosmetics containing microbeads.\n"
    "  – Absorption of smaller nanoplastics (<100 nm) is under investigation."
)

# --- Slide 5: Confirmed Presence in the Human Body ---
slide_layout = prs.slide_layouts
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.shapes.placeholders

title.text = "Confirmed Presence: Microplastics Inside the Human Body"
content.text = (
    "• Internalization is now an empirically confirmed fact, indicating systemic exposure.\n"
    "• Evidence of Systemic Exposure:\n"
    "  – Gastrointestinal Tract: Found in human stool samples.\n"
    "  – Systemic Circulation: Detected in the blood of 77% of healthy volunteers.\n"
    "  – Lungs: MPs, especially fibers, found in human lung tissue.\n"
    "  – Placenta: Detected in human placental tissue, confirming fetal exposure.\n"
    "  – Other Organs: Identified in human liver, kidney, and spleen tissues."
)

# --- Slide 6: Mechanisms of Toxicity ---
# Visual Suggestion: A central graphic of a microplastic particle with two branching arms.
# We will use a two-content layout for this.
slide_layout = prs.slide_layouts # Two Content layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Mechanisms of Toxicity: How Microplastics Cause Harm"

# Left content box for Physical Effects
left_content = slide.placeholders
left_content.text = "1. Physical Effects (The Particle Itself)"
p = left_content.text_frame.add_paragraph()
p.text = "• Causes cellular stress, membrane damage, inflammation, and cell death."
p.level = 1
p = left_content.text_frame.add_paragraph()
p.text = "• In lungs, can induce granulomas and fibrosis."
p.level = 1
p = left_content.text_frame.add_paragraph()
p.text = "• Nanoplastics can cross blood-brain and placental barriers."
p.level = 1

# Right content box for Chemical Effects
right_content = slide.placeholders
right_content.text = "2. Chemical Effects (The Particle as a Carrier)"
p = right_content.text_frame.add_paragraph()
p.text = "• Leaching of Additives: Harmful chemicals like BPA and phthalates (endocrine disruptors) can leach out."
p.level = 1
p = right_content.text_frame.add_paragraph()
p.text = "• Vector for Contaminants: Absorb and transport pollutants (pesticides, heavy metals) into tissues."
p.level = 1

# --- Slide 7: Key Health Effects ---
slide_layout = prs.slide_layouts
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.shapes.placeholders

title.text = "Key Health Effects (Evidence from In Vitro & In Vivo Models)"
content.text = (
    "• Oxidative Stress and Inflammation:\n"
    "  – The most universally reported mechanism of MP toxicity.\n"
    "  – Induces reactive oxygen species (ROS), leading to cellular damage.\n"
    "• Other Major Concerns:\n"
    "  – Metabolic Disruption: Gut microbiota dysbiosis, altered lipid metabolism.\n"
    "  – Immunotoxicity: May alter the function of immune cells.\n"
    "  – Reproductive Toxicity: Reduced fertility and sperm quality in animal models.\n"
    "  – Neurotoxicity: Can cause neuroinflammation and disrupt neurotransmitters."
)

# --- Slide 8: Critical Knowledge Gaps ---
slide_layout = prs.slide_layouts # Title and Content
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders

title.text = "The Path Forward: Critical Knowledge Gaps"

# Visual Suggestion: Four content areas for the gaps. We will use bullet points.
content.text = (
    "1. Analytical Challenges: Lack of standardized and validated methods for identifying and quantifying MPs, especially nanoplastics, in human tissues.\n\n"
    "2. Lack of Epidemiological Data: A severe scarcity of large-scale human studies linking internal MP concentrations to specific diseases.\n\n"
    "3. Poorly Understood Toxicokinetics (ADME): How MPs are Absorbed, Distributed, Metabolized, and Excreted is not well characterized.\n\n"
    "4. Unknown Dose-Response Relationships: The threshold at which MP exposure causes adverse effects in humans is unknown."
)
# Add a picture placeholder - You will need to create/find an image for this.
# For example, a picture of question marks.
# You need to have 'image_gaps.jpg' in the same folder as the script.
try:
    img_path = 'image_gaps.jpg' #<-- CHANGE THIS TO YOUR IMAGE FILE
    # Adjust left, top, width, height as needed
    slide.shapes.add_picture(img_path, Inches(7), Inches(2), width=Inches(5))
except FileNotFoundError:
    print("Warning: 'image_gaps.jpg' not found. Skipping image on Slide 8.")


# --- Slide 9: Conclusion and Recommendations ---
slide_layout = prs.slide_layouts
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.shapes.placeholders

title.text = "Conclusion and Recommendations"
content.text = (
    "• The Reality: The infiltration of microplastics into the human body is a confirmed reality, not speculation.\n"
    "• The Risk: Toxicological evidence shows MPs can cause adverse biological effects, primarily through oxidative stress and inflammation.\n"
    "• Call to Action: The precautionary principle warrants a concerted, global effort:\n"
    "  1. Science: Robust research to close critical knowledge gaps.\n"
    "  2. Policy: Regulatory action to reduce plastic production and improve waste management.\n"
    "  3. Public: Awareness campaigns to drive behavioral change."
)


# --- Slide 10: Thank You & Questions ---
slide_layout = prs.slide_layouts # Section Header layout (can be used for end slide)
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders

title.text = "Thank You"
subtitle.text = "Questions?\n\nContact: Dr. Siddalingaiah H S | hssling@yahoo.com"

# --- Save the presentation ---
file_name = "Microplastics_and_Human_Health.pptx"
prs.save(file_name)
print(f"Presentation saved as {file_name}")

```